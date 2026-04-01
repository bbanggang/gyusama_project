"""
각 주차 README.md 기반 발표용 PPT 생성 스크립트
실행: python3 make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # 짙은 네이비
C_ACCENT   = RGBColor(0x00, 0xC8, 0xFF)   # 시안
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xB0, 0xC4, 0xDE)   # 연한 스틸블루
C_GREEN    = RGBColor(0x00, 0xE5, 0x96)   # 완료 색상
C_CARD     = RGBColor(0x16, 0x2D, 0x45)   # 카드 배경
C_YELLOW   = RGBColor(0xFF, 0xD6, 0x00)   # 강조

W, H = Inches(13.33), Inches(7.5)         # 와이드 16:9


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, l, t, w, h, font_size, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_divider(slide, t, color=C_ACCENT):
    add_rect(slide, Inches(0.5), t, Inches(12.33), Inches(0.04), color)


# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 제작 함수들
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs, week_num, title, subtitle, goal):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # 상단 액센트 바
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_ACCENT)

    # WEEK 배지
    add_rect(slide, Inches(0.6), Inches(1.2), Inches(1.8), Inches(0.55), C_ACCENT)
    add_text(slide, f"WEEK {week_num}", Inches(0.6), Inches(1.2), Inches(1.8), Inches(0.55),
             14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # 제목
    add_text(slide, title, Inches(0.6), Inches(1.9), Inches(11.5), Inches(1.4),
             36, bold=True, color=C_WHITE)

    # 부제목
    add_text(slide, subtitle, Inches(0.6), Inches(3.3), Inches(11.0), Inches(0.6),
             18, color=C_ACCENT)

    add_divider(slide, Inches(4.1))

    # 목표
    add_text(slide, "목표", Inches(0.6), Inches(4.3), Inches(1.5), Inches(0.4),
             13, bold=True, color=C_ACCENT)
    add_text(slide, goal, Inches(0.6), Inches(4.75), Inches(11.5), Inches(1.5),
             14, color=C_LIGHT)

    # 하단 바
    add_rect(slide, Inches(0), Inches(7.3), W, Inches(0.2), C_ACCENT)


def slide_env_check(prs, items):
    """환경 점검 슬라이드 (표 형태)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_ACCENT)

    add_text(slide, "01  환경 점검", Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
             24, bold=True, color=C_ACCENT)
    add_divider(slide, Inches(1.0))

    row_h = Inches(0.52)
    for i, (item, cmd, result) in enumerate(items):
        y = Inches(1.15) + i * row_h
        bg = C_CARD if i % 2 == 0 else C_BG
        add_rect(slide, Inches(0.5), y, Inches(12.3), row_h - Inches(0.04), bg)
        add_text(slide, item,   Inches(0.6),  y + Inches(0.08), Inches(4.0),  row_h, 11, color=C_WHITE)
        add_text(slide, cmd,    Inches(4.7),  y + Inches(0.08), Inches(4.5),  row_h, 10, color=C_LIGHT)
        add_text(slide, result, Inches(9.3),  y + Inches(0.08), Inches(3.3),  row_h, 11,
                 bold=True, color=C_GREEN)

    add_rect(slide, Inches(0), Inches(7.3), W, Inches(0.2), C_ACCENT)


def slide_three_cols(prs, slide_title_text, slide_num, col_data):
    """3컬럼 작업 내용 슬라이드 [(제목, [bullet, ...])]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_ACCENT)

    add_text(slide, f"{slide_num:02d}  {slide_title_text}",
             Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             24, bold=True, color=C_ACCENT)
    add_divider(slide, Inches(1.0))

    col_w = Inches(3.9)
    gap   = Inches(0.27)
    for i, (col_title, bullets) in enumerate(col_data):
        x = Inches(0.5) + i * (col_w + gap)
        # 카드 배경
        add_rect(slide, x, Inches(1.1), col_w, Inches(5.9), C_CARD)
        # 컬럼 제목 배경
        add_rect(slide, x, Inches(1.1), col_w, Inches(0.55), C_ACCENT)
        add_text(slide, col_title, x + Inches(0.1), Inches(1.12), col_w - Inches(0.2), Inches(0.5),
                 13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 불릿
        bullet_text = "\n".join(f"• {b}" for b in bullets)
        add_text(slide, bullet_text, x + Inches(0.15), Inches(1.75),
                 col_w - Inches(0.3), Inches(5.1), 11, color=C_LIGHT)

    add_rect(slide, Inches(0), Inches(7.3), W, Inches(0.2), C_ACCENT)


def slide_completion(prs, checks, notes):
    """완료 기준 + 주의사항 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_ACCENT)

    add_text(slide, "완료 기준 확인 & 주의사항",
             Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             24, bold=True, color=C_ACCENT)
    add_divider(slide, Inches(1.0))

    # 완료 기준 (좌)
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.9), Inches(5.9), C_CARD)
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.9), Inches(0.55), C_GREEN)
    add_text(slide, "✓  완료 기준", Inches(0.6), Inches(1.12), Inches(5.7), Inches(0.5),
             13, bold=True, color=C_BG)
    for i, check in enumerate(checks):
        y = Inches(1.8) + i * Inches(0.62)
        add_rect(slide, Inches(0.7), y, Inches(0.45), Inches(0.42), C_GREEN)
        add_text(slide, "✓", Inches(0.7), y, Inches(0.45), Inches(0.42),
                 12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, check, Inches(1.25), y, Inches(4.9), Inches(0.5), 11, color=C_WHITE)

    # 주의사항 (우)
    add_rect(slide, Inches(6.93), Inches(1.1), Inches(5.9), Inches(5.9), C_CARD)
    add_rect(slide, Inches(6.93), Inches(1.1), Inches(5.9), Inches(0.55), C_YELLOW)
    add_text(slide, "⚠  주의사항", Inches(7.03), Inches(1.12), Inches(5.7), Inches(0.5),
             13, bold=True, color=C_BG)
    for i, note in enumerate(notes):
        y = Inches(1.8) + i * Inches(0.75)
        add_text(slide, f"▸  {note}", Inches(7.1), y, Inches(5.5), Inches(0.7),
                 11, color=C_LIGHT)

    add_rect(slide, Inches(0), Inches(7.3), W, Inches(0.2), C_ACCENT)


def slide_two_cols(prs, title_text, slide_num, left_title, left_bullets,
                   right_title, right_bullets):
    """2컬럼 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_ACCENT)
    add_text(slide, f"{slide_num:02d}  {title_text}",
             Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
             24, bold=True, color=C_ACCENT)
    add_divider(slide, Inches(1.0))

    for i, (col_t, bullets) in enumerate([(left_title, left_bullets),
                                           (right_title, right_bullets)]):
        x = Inches(0.5) + i * Inches(6.68)
        add_rect(slide, x, Inches(1.1), Inches(6.15), Inches(5.9), C_CARD)
        add_rect(slide, x, Inches(1.1), Inches(6.15), Inches(0.55), C_ACCENT)
        add_text(slide, col_t, x + Inches(0.1), Inches(1.12),
                 Inches(5.9), Inches(0.5), 13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt = "\n".join(f"• {b}" for b in bullets)
        add_text(slide, txt, x + Inches(0.15), Inches(1.75),
                 Inches(5.85), Inches(5.1), 11, color=C_LIGHT)

    add_rect(slide, Inches(0), Inches(7.3), W, Inches(0.2), C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# WEEK 1
# ══════════════════════════════════════════════════════════════════════════════
def make_week1():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs, 1,
        "PC 개발 환경 최적화 및 Docker 기초",
        "Ubuntu 24.04 · ROS2 Jazzy · NVIDIA Container Toolkit · Docker",
        "Ubuntu 24.04에서 GPU 가속 Docker 환경을 구축하고,\nROS2가 포함된 팀 전용 베이스 이미지를 만든다.")

    slide_env_check(prs, [
        ("Ubuntu 24.04.4 LTS",           "lsb_release -a",                       "✓ 정상"),
        ("RTX 5070 Ti  Driver 590 / CUDA 13.1",  "nvidia-smi",                   "✓ 정상"),
        ("Docker 29.3.1 + NVIDIA Runtime","docker info | grep Runtime",           "✓ 정상"),
        ("NVIDIA Container Toolkit 1.19.0","nvidia-ctk --version",               "✓ 정상"),
        ("Isaac Sim 5.1.0  (~/isaac_env)","pip show isaacsim | grep Version",    "✓ 정상"),
        ("PyTorch 2.7.0+cu128 / ONNX 1.20.1","pip list | grep -E 'torch|onnx'", "✓ 정상"),
        ("RAM 62GB / 디스크 여유 126GB",  "free -h && df -h /",                   "✓ 충분"),
        ("ROS2 Jazzy (호스트)",           "ls /opt/ros/jazzy/setup.bash",         "✓ 설치"),
    ])

    slide_three_cols(prs, "진행 작업", 2, [
        ("1-2. ROS2 Jazzy 설치", [
            "호스트에서 ros2 명령 직접 사용",
            "ros-jazzy-desktop 전체 설치",
            "TurtleBot3 / Nav2 패키지 포함",
            "ROS_DOMAIN_ID=1 환경변수 설정",
            "TURTLEBOT3_MODEL=burger 설정",
        ]),
        ("1-3. 프로젝트 구조 생성", [
            "docker/ — Dockerfile 관리",
            "src/ — ROS2 노드 소스",
            "isaac_sim/ — 시뮬레이터 스크립트",
            "models/ — AI 모델",
            "data/{synthetic, real, calibration}",
            "configs/ — 설정 파일",
        ]),
        ("1-4. Dockerfile.pc 작성", [
            "Base: nvidia/cuda:12.6.0-ubuntu24.04",
            "ROS2 Jazzy Desktop + TurtleBot3",
            "Nav2 패키지 포함",
            "ROS_DOMAIN_ID / TURTLEBOT3_MODEL 내장",
            "entrypoint.sh 자동 소싱",
            "docker build 빌드 성공",
        ]),
    ])

    slide_completion(prs,
        checks=[
            "docker run --gpus all → RTX 5070 Ti 인식",
            "Dockerfile.pc 빌드 성공 (FINISHED)",
            "컨테이너 내 ros2 topic list 정상 출력",
        ],
        notes=[
            "ROS_DOMAIN_ID는 PC·컨테이너·RPi5 모두 1로 통일",
            "ENV PYTHONPATH=...:$PYTHONPATH 형식 사용 금지\n→ entrypoint.sh에서 export 방식 사용",
            "재부팅 후 QEMU binfmt 초기화됨\n→ ARM64 빌드 전 재등록 필요",
        ]
    )

    prs.save("week1/presentation.pptx")
    print("[✓] week1/presentation.pptx 저장 완료")


# ══════════════════════════════════════════════════════════════════════════════
# WEEK 2
# ══════════════════════════════════════════════════════════════════════════════
def make_week2():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs, 2,
        "Multi-arch 빌드 파이프라인 구축",
        "Docker Buildx · QEMU · ARM64 · x86_64",
        "PC(x86_64)에서 빌드한 코드가 RPi5(ARM64)에서도 실행되도록\n멀티 아키텍처 빌드 환경을 구성한다.")

    slide_two_cols(prs, "핵심 개념", 1,
        "왜 멀티아키텍처가 필요한가?",
        [
            "PC CPU : x86_64 (Intel/AMD)",
            "RPi5 CPU : ARM64 (Cortex-A76)",
            "CPU 아키텍처가 다르면 바이너리 호환 불가",
            "Docker 이미지도 아키텍처별로 다름",
            "Buildx + QEMU → 하나의 명령으로 두 이미지 동시 생성",
            "Docker Hub manifest list → 자동 아키텍처 선택",
        ],
        "QEMU 에뮬레이션이란?",
        [
            "QEMU : 다른 CPU를 소프트웨어로 흉내내는 에뮬레이터",
            "PC에서 ARM64 명령어를 실행 가능하게 함",
            "binfmt_misc 커널 기능으로 투명하게 동작",
            "단점 : 에뮬레이션이므로 네이티브보다 느림",
            "onnxruntime ARM64 pip 빌드 : 약 18분 소요",
            "재부팅 후 binfmt 초기화 → 재등록 필요",
        ]
    )

    slide_three_cols(prs, "진행 작업", 2, [
        ("2-1. QEMU 등록", [
            "multiarch/qemu-user-static 이미지 실행",
            "aarch64-static binfmt 등록 확인",
            "--reset -p yes 플래그로 전체 재등록",
            "ARM64 에뮬레이션 활성화",
        ]),
        ("2-2. Buildx 빌더 생성", [
            "gyusama-builder 전용 빌더 생성",
            "docker-container 드라이버 사용",
            "지원 플랫폼: amd64 + arm64 확인",
            "default 빌더와 분리하여 관리",
        ]),
        ("2-3. Dockerfile.rpi5 작성", [
            "Base: ros:jazzy-ros-base (공식 멀티아치)",
            "TurtleBot3 + Dynamixel SDK 패키지",
            "onnxruntime CPU 버전 설치",
            "QEMU 에뮬레이션 BrokenPipeError 우회",
            "--no-install-recommends 필수",
        ]),
    ])

    slide_completion(prs,
        checks=[
            "docker buildx ls → gyusama-builder amd64+arm64 확인",
            "x86_64 + ARM64 동시 빌드 FINISHED",
            "두 아키텍처 모두 빌드 성공",
        ],
        notes=[
            "QEMU py3compile BrokenPipeError\n→ ros:jazzy-ros-base 공식 이미지 사용으로 해결",
            "--no-install-recommends 미사용 시\n불필요한 패키지(cppcheck, flake8)가 설치되어 이미지 비대화",
            "QEMU binfmt는 재부팅 후 초기화\n→ 영구 등록: binfmt-support 패키지 사용",
            "ARM64 onnxruntime pip 빌드 시간 약 18분\n→ Dockerfile 레이어 구조로 캐시 활용 권장",
        ]
    )

    prs.save("week2/presentation.pptx")
    print("[✓] week2/presentation.pptx 저장 완료")


# ══════════════════════════════════════════════════════════════════════════════
# WEEK 3
# ══════════════════════════════════════════════════════════════════════════════
def make_week3():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs, 3,
        "원격 배포 및 컨테이너 오케스트레이션",
        "Docker Hub · docker-compose · FastDDS · ROS2 토픽 통신",
        "Docker Hub로 이미지를 관리하고, docker-compose로 ROS2 노드를\n한 번에 배포하여 PC ↔ RPi5 통신을 확인한다.")

    slide_two_cols(prs, "Docker Hub Push", 1,
        "멀티아키텍처 이미지 Push",
        [
            "QEMU 재등록 (재부팅 후 초기화)",
            "gyusama-builder 재생성 및 ARM64 지원 확인",
            "docker login -u bbanggang",
            "buildx build --platform linux/amd64,linux/arm64",
            "--push 플래그로 Hub에 직접 업로드",
            "Digest: sha256:2453c... 확인",
            "RPi5 pull 시 ARM64 자동 선택",
        ],
        "새 RPi5 보드 환경 구성",
        [
            "기존 보드: 네이티브 ROS2 충돌 위험",
            "새 보드에 Ubuntu 클린 설치",
            "SSH 연결 확인 (192.168.0.155)",
            "Docker 29.3.1 설치",
            "Docker Compose v5.1.1 설치",
            "sudo usermod -aG docker → newgrp docker",
        ]
    )

    slide_three_cols(prs, "docker-compose 서비스 통합", 2, [
        ("ros2-core", [
            "ROS2 데몬 + 통신 기반",
            "network_mode: host",
            "ROS_DOMAIN_ID=1",
            "TURTLEBOT3_MODEL=burger",
            "restart: unless-stopped",
        ]),
        ("control-node", [
            "Dynamixel 모터 제어",
            "privileged: true",
            "/dev/ttyUSB0 디바이스 마운트",
            "LDS_MODEL=LDS-02 필수",
            "ros2 launch turtlebot3_bringup",
        ]),
        ("설계 핵심 결정", [
            "network_mode: host\n→ FastDDS 멀티캐스트 통과",
            "privileged: true\n→ /dev/ttyUSB0, /dev/video0",
            "LDS_MODEL 누락 시 KeyError 크래시",
            "하드웨어 미연결 시 sleep infinity",
            "ROS_DOMAIN_ID 전 노드 통일 필수",
        ]),
    ])

    slide_completion(prs,
        checks=[
            "Docker Hub bbanggang/turtlebot3-ros2:v1.0 push 성공",
            "RPi5에서 docker pull → ARM64 자동 선택",
            "docker compose up → ros2-core + control-node 기동",
            "PC에서 ros2 topic echo → hello_from_rpi5 수신",
        ],
        notes=[
            "network_mode: host 필수\n→ bridge 네트워크는 FastDDS 멀티캐스트 차단",
            "LDS_MODEL 환경변수 누락 시\nrobot.launch.py KeyError로 컨테이너 재시작 루프",
            "기존 RPi5 사용 시 네이티브 ROS2 데몬 충돌\n→ ros2 daemon stop 후 테스트",
            "ROS_DOMAIN_ID 불일치 시\n동일 네트워크에서도 토픽 수신 불가",
        ]
    )

    prs.save("week3/presentation.pptx")
    print("[✓] week3/presentation.pptx 저장 완료")


if __name__ == "__main__":
    import os
    os.chdir("/home/linux/gyusama-project/weeks_work")
    make_week1()
    make_week2()
    make_week3()
    print("\n모든 PPT 생성 완료!")
